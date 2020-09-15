import os
import pptx
from pptx.util import Inches


def modify_path(input_path, suffix="_edit"):
    """
    This function modifies an file path by inserting the given suffix between the base file name
    and the file extension. For example, C:/dog/cat.txt becomes C:/dog/cat_edit.txt, for a
    suffix '_edit'

    :param input_path: the path string to modify
    :param suffix: the string to insert into the path
    :return: the modified path string
    """
    head, tail = os.path.split(input_path)
    basename, extension = os.path.splitext(tail)
    return os.path.join(head, basename + suffix + extension)


def find_images(folder_path, extensions=('.png', '.jpg', '.tif')):
    """
    This function searches a folder for files that have one of the given file extensions.

    :param folder_path: a string path to the images folder
    :param extensions: a list of image extensions strings like ['.png', '.jpg']
    :return: a list of image file names strings
    """
    image_full_paths = []
    for filename in os.listdir(folder_path):
        basename, extension = os.path.splitext(filename)
        if extension.lower() in extensions:
            image_full_paths.append(os.path.join(folder_path, filename))
    return image_full_paths


def create_blank_prs(slide_height_inch=7.5, slide_width_inch=13.333):
    """
    This function creates a new Presentation object with specified slide dimensions.

    :param slide_height_inch: the height of the slide in inches
    :param slide_width_inch: the width of the slide in inches
    :return: a Presentation object
    """
    # Create the Blank Presentation
    prs = pptx.Presentation()
    prs.slide_height = Inches(slide_height_inch)
    prs.slide_width = Inches(slide_width_inch)
    return prs


def add_image_slide(prs, img_path):
    """
    This function creates a new blank slide and adds a image to it. The image is rescaled to
    fit the width of the slide, while retaining its aspect ratio.

    :param prs: the Presentation object
    :param img_path: a string containing the path to the image file
    :return: a Presentation object
    """
    # Add a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # Add the image, resizing it to fit slide width, while maintaining aspect ratio
    new_slide.shapes.add_picture(img_path, left=Inches(0), top=Inches(0), width=prs.slide_width)
    return prs


def create_slides(img_folder_path, pptx_path, slide_height_inch=7.5, slide_width_inch=13.333,
                  extensions=('.png', '.jpg', '.tif')):
    """
    This function can create image-based slides and add them to a power point presentation in one of
    three modes.
    1) Create: If the given presentation path does not exist, a new presentation is created
    2) Copy and Append: If the given presentation path does exist, a copy is made, new slides are appended, and
    the presentation is saved with the original filename plus an '_edit' suffix.

    :param img_folder_path: a string path to the images folder
    :param pptx_path: a string path for the PowerPoint file
    :param slide_height_inch: the height of the slide in inches
    :param slide_width_inch: the width of the slide in inches
    :param extensions: a list of image extensions strings like ['.png', '.jpg']
    :return: None
    """
    # Validate the image folder path
    if not os.path.isdir(img_folder_path):
        print('Cannot proceed because the image folder path is invalid')
        return

    # Load existing presentation
    found_existing = True
    try:
        prs = pptx.Presentation(pptx_path)
        print("Found existing file, will copy and append with new slides")

    # Or Create a blank presentation
    except pptx.exc.PackageNotFoundError:
        found_existing = False
        prs = create_blank_prs(slide_height_inch, slide_width_inch)
        print("No existing file found, will create one with new slides")

    # Add all of the image slides
    for img_path in find_images(img_folder_path, extensions):
        add_image_slide(prs, img_path)
        print("Added image slide containing image: '{0}'".format(img_path))

    # Save the presentation
    if found_existing:
        pptx_path = modify_path(pptx_path)
    prs.save(pptx_path)
    print("\nSaved file to: '{0}'".format(pptx_path))
